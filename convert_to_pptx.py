#!/usr/bin/env python3
"""
Script pour convertir PRESENTATION_DIVALEO.md en PowerPoint
Nécessite: pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    import re
except ImportError:
    print("❌ Erreur: python-pptx n'est pas installé")
    print("📦 Installation: pip install python-pptx")
    exit(1)

def parse_markdown_slides(md_file):
    """Parse le fichier Markdown et extrait les slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Diviser par les séparateurs de slides
    slides = re.split(r'^## SLIDE \d+', content, flags=re.MULTILINE)
    
    parsed_slides = []
    for i, slide_content in enumerate(slides[1:], 1):  # Ignorer le premier élément (en-tête)
        lines = slide_content.strip().split('\n')
        
        # Extraire le titre (première ligne après ":")
        title = ""
        content_lines = []
        
        for line in lines:
            line = line.strip()
            if not line or line.startswith('---'):
                continue
            
            # Titre (ligne qui commence par ** ou #)
            if line.startswith('**') and not title:
                title = line.replace('**', '').replace('*', '').strip()
            elif line.startswith('###') and not title:
                title = line.replace('###', '').strip()
            elif ':' in line and not title and len(line) < 100:
                title = line.split(':')[1].strip() if ':' in line else line
            else:
                # Contenu
                clean_line = line.replace('**', '').replace('*', '').replace('✅', '✓').replace('🔐', '').replace('📊', '').replace('👥', '').replace('🏗️', '').replace('📝', '').replace('📦', '').replace('📋', '').replace('🔒', '').replace('🛡️', '').replace('🔄', '').replace('📈', '').replace('📉', '').replace('🎯', '').replace('⚡', '').replace('🚀', '').replace('🌐', '').replace('📱', '').replace('🎨', '').replace('🖱️', '').replace('🔔', '').replace('👑', '').replace('👁️', '').replace('👔', '').replace('➕', '').replace('✏️', '').replace('📄', '').replace('📎', '').replace('🔄', '').replace('📊', '').replace('📈', '').replace('📉', '').replace('🎯', '').replace('🔐', '').replace('🔑', '').replace('🛡️', '').replace('🔒', '').replace('📝', '').replace('🔗', '').replace('✅', '✓').replace('🌍', '').replace('☁️', '').replace('🤖', '').replace('📧', '').replace('💾', '').replace('👥', '').replace('📊', '').replace('📈', '').replace('⏱️', '').replace('📚', '').replace('🎥', '').replace('📖', '').replace('📧', '').replace('📱', '').replace('🌐', '').strip()
                if clean_line:
                    content_lines.append(clean_line)
        
        if not title and content_lines:
            title = content_lines[0]
            content_lines = content_lines[1:]
        
        parsed_slides.append({
            'number': i,
            'title': title or f'Slide {i}',
            'content': '\n'.join(content_lines)
        })
    
    return parsed_slides

def create_powerpoint(slides, output_file):
    """Crée une présentation PowerPoint à partir des slides parsés"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        # Utiliser le layout "Title and Content"
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Titre
        title_shape = slide.shapes.title
        title_shape.text = slide_data['title']
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Contenu
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        # Diviser le contenu en paragraphes
        content_lines = slide_data['content'].split('\n')
        
        for i, line in enumerate(content_lines):
            if line.strip():
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = line.strip()
                p.font.size = Pt(18)
                p.level = 0
                
                # Détecter les listes à puces
                if line.strip().startswith('-') or line.strip().startswith('•'):
                    p.level = 1
                    p.text = line.strip()[1:].strip()
                    p.font.size = Pt(16)
    
    prs.save(output_file)
    print(f"✅ PowerPoint créé avec succès : {output_file}")
    print(f"📊 Nombre de slides : {len(slides)}")

if __name__ == "__main__":
    input_file = "PRESENTATION_DIVALEO.md"
    output_file = "PRESENTATION_DIVALEO.pptx"
    
    print("🔄 Conversion de Markdown vers PowerPoint...")
    print(f"📄 Fichier source : {input_file}")
    print(f"📊 Fichier de sortie : {output_file}\n")
    
    try:
        slides = parse_markdown_slides(input_file)
        print(f"✅ {len(slides)} slides détectés\n")
        
        create_powerpoint(slides, output_file)
        
        print("\n🎉 Conversion terminée avec succès !")
        print(f"📥 Vous pouvez maintenant ouvrir : {output_file}")
        
    except FileNotFoundError:
        print(f"❌ Erreur : Le fichier {input_file} est introuvable")
        print("💡 Assurez-vous que le fichier est dans le même dossier que ce script")
    except Exception as e:
        print(f"❌ Erreur lors de la conversion : {e}")
        print("\n💡 Alternative : Utilisez la méthode manuelle (Option 1)")

